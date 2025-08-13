import streamlit as st
import google.generativeai as genai
import PyPDF2
from pptx import Presentation
import io
import base64
from PIL import Image
import tempfile
import os
import json
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak, Image as RLImage
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY
import markdown
import re
from functools import lru_cache
from contextlib import contextmanager
import logging
from typing import Dict, Optional, Tuple, Any, List
import requests
from openai import OpenAI
import time

# Import custom modules
from auth import check_authentication, show_logout_option
from database import log_ui_interaction, update_ui_interaction
from prompts import (
    DOCUMENT_ANALYSIS_PROMPT,
    VIDEO_SCRIPT_PROMPT,
    CONCLUSION_PROMPT,
    ASSESSMENT_PROMPT,
    DOUBT_RESOLUTION_PROMPT,
    PDF_CONTENT_PROMPT,
    TRANSLATION_PROMPT,
    IMAGE_PROMPT_GENERATION_PROMPT 
)

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Constants
SUPPORTED_FILE_TYPES = ['pdf', 'pptx', 'ppt']
MAX_FILE_SIZE = 10 * 1024 * 1024  # 10MB
GEMINI_MODEL = 'gemini-2.0-flash'
OPENAI_API_KEY = "openai-api-key"
SYNTHESIA_API_KEY = "synthesia-api-key" 
SYNTHESIA_BASE_URL = "https://api.synthesia.io/v2"

# Configure page
st.set_page_config(
    page_title="Student Document Analyzer",
    page_icon="📚",
    layout="wide",
    initial_sidebar_state="expanded"
)

class SessionManager:
    """Manages Streamlit session state efficiently"""
    
    DEFAULT_VALUES = {
        'analysis_completed': False,
        'document_text': "",
        'analysis_result': "",
        'assessment_questions': None,
        'show_doubt_session': False,
        'show_assessment': False,
        'show_video_script': False,
        'show_video_generation': False,  
        'show_conclusion': False,
        'show_personalized_pdf': False,
        'video_script': "",
        'video_scripts_by_level': {},  
        'generated_videos': {}, 
        'conclusion_content': "",
        'quiz_performance': None,
        'selected_language': "English",
        'language_code': "en",
        'user_email': ""
    }
    
    @staticmethod
    def initialize():
        """Initialize session state with default values"""
        for key, value in SessionManager.DEFAULT_VALUES.items():
            if key not in st.session_state:
                st.session_state[key] = value
    
    @staticmethod
    def reset_analysis():
        """Reset analysis-related session state"""
        reset_keys = [
            'analysis_completed', 'document_text', 'analysis_result', 
            'assessment_questions', 'video_script', 'video_scripts_by_level',
            'generated_videos', 'conclusion_content', 'quiz_performance'
        ]
        for key in reset_keys:
            if key in SessionManager.DEFAULT_VALUES:
                st.session_state[key] = SessionManager.DEFAULT_VALUES[key]
        
        # Reset all show flags
        show_keys = [k for k in SessionManager.DEFAULT_VALUES.keys() if k.startswith('show_')]
        for key in show_keys:
            st.session_state[key] = False
    
    @staticmethod
    def set_view(view_name: str):
        """Set active view and reset others"""
        show_keys = [k for k in SessionManager.DEFAULT_VALUES.keys() if k.startswith('show_')]
        for key in show_keys:
            st.session_state[key] = (key == f'show_{view_name}')

class SynthesiaVideoGenerator:
    """Manages Synthesia API integration for video generation"""
    
    def __init__(self):
        self.api_key = SYNTHESIA_API_KEY
        self.base_url = SYNTHESIA_BASE_URL
        self.headers = {
            'Authorization': f'Bearer {self.api_key}',
            'Content-Type': 'application/json'
        }
        
        # Available avatars - you can get these from Synthesia API
        self.avatars = {
            "easy": "anna_costume1_cameraA",  # Friendly, approachable avatar
            "medium": "james_costume1_cameraA",  # Professional avatar
            "hard": "sophia_costume1_cameraA"  # Academic-looking avatar
        }
        
        # Voice settings for different levels
        self.voices = {
            "easy": "en-US-JennyNeural",
            "medium": "en-US-AriaNeural", 
            "hard": "en-US-GuyNeural"
        }
    
    def get_available_avatars(self) -> List[Dict]:
        """Get list of available avatars from Synthesia API"""
        try:
            response = requests.get(
                f"{self.base_url}/avatars",
                headers=self.headers
            )
            if response.status_code == 200:
                return response.json().get('avatars', [])
            else:
                logger.error(f"Failed to get avatars: {response.status_code}")
                return []
        except Exception as e:
            logger.error(f"Error getting avatars: {e}")
            return []
    
    def create_video(self, script: str, difficulty_level: str, title: str) -> Optional[Dict]:
        """Create a video using Synthesia API"""
        try:
            # Prepare the video creation payload
            video_data = {
                "title": f"{title} - {difficulty_level.title()} Level",
                "description": f"Educational video generated for {difficulty_level} difficulty level",
                "input": [
                    {
                        "avatarSettings": {
                            "avatar": self.avatars.get(difficulty_level, self.avatars["medium"]),
                            "horizontalAlign": "center",
                            "scale": 1.0,
                            "style": "rectangular"
                        },
                        "background": {
                            "type": "color",
                            "value": "#f0f8ff" if difficulty_level == "easy" else 
                                    "#fff5ee" if difficulty_level == "medium" else "#f5f5dc"
                        },
                        "scriptText": script,
                        "voiceSettings": {
                            "voice": self.voices.get(difficulty_level, self.voices["medium"]),
                            "speed": 1.0 if difficulty_level != "hard" else 0.9,
                            "pitch": 1.0
                        }
                    }
                ],
                "test": False,  
                "visibility": "private",
                "aspectRatio": "16:9"
            }
            
            # Make API request to create video
            response = requests.post(
                f"{self.base_url}/videos",
                headers=self.headers,
                json=video_data
            )
            
            if response.status_code == 201:
                video_info = response.json()
                return {
                    'video_id': video_info['id'],
                    'status': video_info['status'],
                    'title': video_info['title'],
                    'created_at': video_info['createdAt'],
                    'difficulty_level': difficulty_level
                }
            else:
                logger.error(f"Failed to create video: {response.status_code} - {response.text}")
                return None
                
        except Exception as e:
            logger.error(f"Error creating video with Synthesia: {e}")
            return None
    
    def get_video_status(self, video_id: str) -> Optional[Dict]:
        """Check video generation status"""
        try:
            response = requests.get(
                f"{self.base_url}/videos/{video_id}",
                headers=self.headers
            )
            
            if response.status_code == 200:
                return response.json()
            else:
                logger.error(f"Failed to get video status: {response.status_code}")
                return None
                
        except Exception as e:
            logger.error(f"Error getting video status: {e}")
            return None
    
    def get_video_download_url(self, video_id: str) -> Optional[str]:
        """Get download URL for completed video"""
        try:
            video_info = self.get_video_status(video_id)
            if video_info and video_info.get('status') == 'complete':
                return video_info.get('download')
            return None
        except Exception as e:
            logger.error(f"Error getting download URL: {e}")
            return None

class GeminiManager:
    """Manages Gemini API interactions"""
    
    def __init__(self):
        self._model = None
        self.configure()
    
    def configure(self) -> bool:
        """Configure Gemini API"""
        try:
            api_key = "AIzaSyDsxNln7Fa9aZKlnCYXZKoS5jc_r-Uu5L0"
            genai.configure(api_key=api_key)
            return True
        except Exception as e:
            logger.error(f"Failed to configure Gemini API: {e}")
            return False
    
    @property
    def model(self):
        """Get or create Gemini model instance"""
        if self._model is None:
            self._model = genai.GenerativeModel(GEMINI_MODEL)
        return self._model
    
    def generate_content(self, prompt: str) -> Optional[str]:
        """Generate content using Gemini API with error handling"""
        try:
            response = self.model.generate_content(prompt)
            return response.text
        except Exception as e:
            logger.error(f"Gemini API error: {e}")
            st.error(f"Error with Gemini API: {str(e)}")
            return None

class ImageGenerator:
    """Manages DALL-E image generation"""
    
    def __init__(self):
        self.client = OpenAI(api_key=OPENAI_API_KEY)
        self.gemini = get_gemini_manager()
    
    def generate_image_prompt(self, text_content: str, context: str = "") -> Optional[str]:
        """Generate image prompt using Gemini based on content"""
        try:
            prompt = IMAGE_PROMPT_GENERATION_PROMPT.format(
                content=text_content,
                context=context
            )
            return self.gemini.generate_content(prompt)
        except Exception as e:
            logger.error(f"Error generating image prompt: {e}")
            return None
    
    def generate_image(self, prompt: str, size: str = "1024x1024") -> Optional[bytes]:
        """Generate image using DALL-E API"""
        try:
            response = self.client.images.generate(
                model="dall-e-3",
                prompt=prompt,
                size=size,
                quality="standard",
                n=1,
            )
            
            image_url = response.data[0].url
            
            # Download the image
            img_response = requests.get(image_url)
            if img_response.status_code == 200:
                return img_response.content
            else:
                logger.error(f"Failed to download image: {img_response.status_code}")
                return None
                
        except Exception as e:
            logger.error(f"DALL-E API error: {e}")
            st.error(f"Error generating image: {str(e)}")
            return None
    
    def should_generate_image(self, text_content: str) -> bool:
        """Determine if an image should be generated for the given content"""
        # Keywords that suggest visual content would be helpful
        visual_keywords = [
            'diagram', 'chart', 'graph', 'figure', 'illustration', 'structure',
            'process', 'workflow', 'concept', 'model', 'architecture', 'system',
            'comparison', 'timeline', 'lifecycle', 'hierarchy', 'relationship',
            'mechanism', 'phenomenon', 'experiment', 'observation', 'visualization'
        ]
        
        text_lower = text_content.lower()
        return any(keyword in text_lower for keyword in visual_keywords)

# Global instances
@st.cache_resource
def get_gemini_manager():
    """Get cached Gemini manager instance"""
    return GeminiManager()

@st.cache_resource
def get_image_generator():
    """Get cached image generator instance"""
    return ImageGenerator()

@st.cache_resource
def get_synthesia_generator():
    """Get cached Synthesia video generator instance"""
    return SynthesiaVideoGenerator()

class LanguageManager:
    """Manages language support and translations"""
    
    SUPPORTED_LANGUAGES = {
        "English": "en",
        "Telugu": "te", 
        "Kannada": "kn",
        "Hindi": "hi"
    }
    
    UI_TRANSLATIONS = {
        "en": {
            "upload_header": "📄 Upload Your Document",
            "analysis_button": "🔍 Generate Multi-Level Analysis",
            "doubt_session": "🤔 Doubt Session",
            "assessment": "📝 Assessment",
            "video_script": "🎬 Video Script",
            "video_generation": "🎥 Generate Videos",  
            "conclusion": "🎯 Conclusion",
            "personalized_pdf": "📄 Personalized Course PDF",
            "view_analysis": "📖 View Multi-Level Analysis",
            "new_analysis": "📄 New Analysis",
            "easy_level": "⭐⭐ Easy Level (2.5/5)",
            "medium_level": "⭐⭐⭐⭐ Medium Level (3.5/5)",
            "hard_level": "⭐⭐⭐⭐⭐ Hard Level (4+/5)"
        },
        "te": {
            "upload_header": "📄 మీ పత్రాన్ni అప్‌లోడ్ చేయండి",
            "analysis_button": "🔍 బహుళ-స్థాయి విశ్లేషణ రూపొందించండి",
            "doubt_session": "🤔 సందేహాల సెషన్",
            "assessment": "📝 మూల్యాంకనం",
            "video_script": "🎬 వీడియో స్క్రిప్ట్",
            "video_generation": "🎥 వీడియోలు రూపొందించండి",
            "conclusion": "🎯 ముగింపు",
            "personalized_pdf": "📄 వ్యక్తిగతీకరించిన కోర్స్ PDF",
            "view_analysis": "📖 బహుళ-స్థాయి విశ్లేషణ చూడండి",
            "new_analysis": "📄 కొత్త విశ్లేషణ",
            "easy_level": "⭐⭐ సులభ స్థాయి (2.5/5)",
            "medium_level": "⭐⭐⭐⭐ మధ్యమ స్థాయి (3.5/5)",
            "hard_level": "⭐⭐⭐⭐⭐ కష్ట స్థాయి (4+/5)"
        },
        "hi": {
            "upload_header": "📄 अपना दस्तावेज़ अपलोड करें",
            "analysis_button": "🔍 बहु-स्तरीय विश्लेषण बनाएं",
            "doubt_session": "🤔 संदेह सत्र",
            "assessment": "📝 मूल्यांकन",
            "video_script": "🎬 वीडियो स्क्रिप्ट",
            "video_generation": "🎥 वीडियो बनाएं",
            "conclusion": "🎯 निष्कर्ष",
            "personalized_pdf": "📄 व्यक्तिगत कोर्स PDF",
            "view_analysis": "📖 बहु-स्तरीय विश्लेषण देखें",
            "new_analysis": "📄 नया विश्लेषण",
            "easy_level": "⭐⭐ आसान स्तर (2.5/5)",
            "medium_level": "⭐⭐⭐⭐ मध्यम स्तर (3.5/5)",
            "hard_level": "⭐⭐⭐⭐⭐ कठिन स्तर (4+/5)"
        },
        "kn": {
            "upload_header": "📄 ನಿಮ್ಮ ದಾಖಲೆಯನ್ನು ಅಪ್‌ಲೋಡ್ ಮಾಡಿ",
            "analysis_button": "🔍 ಬಹು-ಹಂತದ ವಿಶ್ಲೇಷಣೆಯನ್ನು ರಚಿಸಿ",
            "doubt_session": "🤔 ಸಂದೇಹ ಸೆಷನ್",
            "assessment": "📝 ಮೌಲ್ಯಮಾಪನ",
            "video_script": "🎬 ವೀಡಿಯೋ ಸ್ಕ್ರಿಪ್ಟ್",
            "video_generation": "🎥 ವೀಡಿಯೋಗಳನ್ನು ರಚಿಸಿ",
            "conclusion": "🎯 ಸಮಾಪನೆ",
            "personalized_pdf": "📄 ವೈಯಕ್ತಿಕಗೊಳಿಸಿದ ಕೋರ್ಸ್ PDF",
            "view_analysis": "📖 ಬಹು-ಹಂತದ ವಿಶ್ಲೇಷಣೆಯನ್ನು ವೀಕ್ಷಿಸಿ",
            "new_analysis": "📄 ಹೊಸ ವಿಶ್ಲೇಷಣೆ",
            "easy_level": "⭐⭐ ಸುಲಭ ಹಂತ (2.5/5)",
            "medium_level": "⭐⭐⭐⭐ ಮಧ್ಯಮ ಹಂತ (3.5/5)",
            "hard_level": "⭐⭐⭐⭐⭐ ಕಠಿಣ ಹಂತ (4+/5)"
        }
    }
    
    @staticmethod
    def get_text(key: str) -> str:
        """Get translated text for UI elements"""
        language_code = st.session_state.get('language_code', 'en')
        translations = LanguageManager.UI_TRANSLATIONS
        return translations[language_code].get(key, translations["en"][key])
    
    @staticmethod
    def translate_content(text: str, target_language: str) -> str:
        """Translate text to target language using Gemini"""
        if target_language == "en" or not text.strip():
            return text
        
        try:
            gemini = get_gemini_manager()
            language_names = {"te": "Telugu", "kn": "Kannada", "hi": "Hindi"}
            
            prompt = TRANSLATION_PROMPT.format(
                target_language=language_names[target_language],
                text=text
            )
            
            return gemini.generate_content(prompt) or text
        except Exception as e:
            logger.error(f"Translation error: {e}")
            return text

class DocumentProcessor:
    """Handles document processing operations"""
    
    @staticmethod
    @contextmanager
    def temporary_file(uploaded_file, suffix='.pptx'):
        """Context manager for temporary file handling"""
        tmp_file = None
        try:
            tmp_file = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)
            tmp_file.write(uploaded_file.read())
            tmp_file.close()
            yield tmp_file.name
        finally:
            if tmp_file and os.path.exists(tmp_file.name):
                os.unlink(tmp_file.name)
    
    @staticmethod
    def extract_text_from_pdf(uploaded_file) -> Optional[str]:
        """Extract text from PDF file with error handling"""
        try:
            pdf_reader = PyPDF2.PdfReader(uploaded_file)
            text_parts = []
            
            for page_num, page in enumerate(pdf_reader.pages, 1):
                try:
                    page_text = page.extract_text()
                    if page_text.strip():
                        text_parts.append(page_text)
                except Exception as e:
                    logger.warning(f"Failed to extract text from page {page_num}: {e}")
                    continue
            
            return "\n".join(text_parts) if text_parts else None
            
        except Exception as e:
            logger.error(f"PDF processing error: {e}")
            st.error(f"Error reading PDF: {str(e)}")
            return None
    
    @staticmethod
    def extract_text_from_ppt(uploaded_file) -> Optional[str]:
        """Extract text from PowerPoint file with error handling"""
        try:
            with DocumentProcessor.temporary_file(uploaded_file, '.pptx') as tmp_path:
                prs = Presentation(tmp_path)
                text_parts = []
                
                for slide_num, slide in enumerate(prs.slides, 1):
                    slide_text = f"Slide {slide_num}:\n"
                    slide_content = []
                    
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_content.append(shape.text.strip())
                    
                    if slide_content:
                        slide_text += "\n".join(slide_content) + "\n\n"
                        text_parts.append(slide_text)
                
                return "\n".join(text_parts) if text_parts else None
                
        except Exception as e:
            logger.error(f"PowerPoint processing error: {e}")
            st.error(f"Error reading PowerPoint: {str(e)}")
            return None
    
    @staticmethod
    def validate_file(uploaded_file) -> Tuple[bool, str]:
        """Validate uploaded file"""
        if uploaded_file.size > MAX_FILE_SIZE:
            return False, f"File size ({uploaded_file.size:,} bytes) exceeds maximum allowed size ({MAX_FILE_SIZE:,} bytes)"
        
        file_type = uploaded_file.type.lower()
        if not any(ft in file_type for ft in ['pdf', 'powerpoint', 'presentation']):
            return False, "Unsupported file type. Please upload PDF or PowerPoint files only."
        
        return True, "Valid file"

class ContentGenerator:
    """Handles AI content generation"""
    
    def __init__(self):
        self.gemini = get_gemini_manager()
        self.image_gen = get_image_generator()
    
    def analyze_document(self, text: str, file_type: str) -> Optional[str]:
        """Analyze document with Gemini API"""
        prompt = DOCUMENT_ANALYSIS_PROMPT.format(file_type=file_type, text=text)
        return self.gemini.generate_content(prompt)
    
    def generate_video_script(self, document_text: str) -> Optional[str]:
        """Generate engaging video script"""
        prompt = VIDEO_SCRIPT_PROMPT.format(document_text=document_text)
        return self.gemini.generate_content(prompt)
    
    def generate_multi_level_video_scripts(self, document_text: str) -> Optional[Dict[str, str]]:
        """Generate three video scripts for different difficulty levels"""
        scripts = {}
        difficulty_levels = {
            "easy": "Create a beginner-friendly video script that explains concepts in simple terms with lots of examples and analogies. Target audience: beginners or students new to the topic.",
            "medium": "Create an intermediate-level video script that covers the main concepts with moderate detail and some technical terms. Target audience: students with some background knowledge.",
            "hard": "Create an advanced-level video script that goes deep into technical details, includes advanced concepts, and uses professional terminology. Target audience: advanced students or professionals."
        }
        
        for level, instruction in difficulty_levels.items():
            prompt = f"""
            {instruction}
            
            Document Content: {document_text}
            
            Please create an engaging educational video script that:
            1. Has a clear introduction, body, and conclusion
            2. Includes natural speaking transitions
            3. Is appropriate for the {level} difficulty level
            4. Is between 300-500 words for optimal video length
            5. Includes engaging elements to keep viewers interested
            
            Format the script as natural speaking text that an AI avatar can deliver effectively.
            """
            
            try:
                script = self.gemini.generate_content(prompt)
                if script:
                    scripts[level] = script
            except Exception as e:
                logger.error(f"Error generating {level} script: {e}")
        
        return scripts if scripts else None
    
    def generate_conclusion(self, document_text: str) -> Optional[str]:
        """Generate comprehensive conclusion"""
        prompt = CONCLUSION_PROMPT.format(document_text=document_text)
        return self.gemini.generate_content(prompt)
    
    def generate_assessment(self, document_text: str) -> Optional[str]:
        """Generate assessment questions"""
        prompt = ASSESSMENT_PROMPT.format(document_text=document_text)
        return self.gemini.generate_content(prompt)
    
    def answer_doubt(self, question: str, document_text: str) -> Optional[str]:
        """Answer student's doubt"""
        prompt = DOUBT_RESOLUTION_PROMPT.format(
            document_text=document_text,
            question=question
        )
        return self.gemini.generate_content(prompt)
    
    def generate_pdf_content(self, document_text: str, quiz_performance=None) -> Optional[Dict]:
        """Generate content for personalized course PDF with image suggestions"""
        performance_context = ""
        if quiz_performance:
            performance_context = f"Student's quiz performance: {quiz_performance}. Tailor the content difficulty and focus areas based on this performance."
        
        prompt = PDF_CONTENT_PROMPT.format(
            performance_context=performance_context,
            document_text=document_text
        )
        
        content = self.gemini.generate_content(prompt)
        if not content:
            return None
        
        # Parse content and identify sections that need images
        sections = self._parse_content_sections(content)
        image_sections = []
        
        for i, section in enumerate(sections):
            if self.image_gen.should_generate_image(section['content']):
                image_prompt = self.image_gen.generate_image_prompt(
                    section['content'], 
                    f"Educational illustration for: {section['title']}"
                )
                if image_prompt:
                    image_sections.append({
                        'section_index': i,
                        'image_prompt': image_prompt,
                        'title': section['title']
                    })
        
        return {
            'content': content,
            'sections': sections,
            'image_sections': image_sections
        }
    
    def _parse_content_sections(self, content: str) -> List[Dict]:
        """Parse content into sections for image generation"""
        sections = []
        lines = content.split('\n')
        current_section = {'title': '', 'content': ''}
        
        for line in lines:
            line = line.strip()
            if line.startswith('# ') or line.startswith('## '):
                if current_section['content']:
                    sections.append(current_section.copy())
                current_section = {
                    'title': line.lstrip('#').strip(),
                    'content': line + '\n'
                }
            else:
                current_section['content'] += line + '\n'
        
        if current_section['content']:
            sections.append(current_section)
        
        return sections

class PDFGenerator:
    """Handles PDF generation operations with image support"""
    
    def __init__(self):
        self.image_gen = get_image_generator()
    
    def create_pdf_from_content(self, content_data: Dict, filename: str = "personalized_course_guide.pdf") -> Optional[bytes]:
        """Create a PDF from the generated content with images"""
        try:
            buffer = io.BytesIO()
            doc = SimpleDocTemplate(buffer, pagesize=letter)
            styles = getSampleStyleSheet()

            # Define custom styles
            title_style = ParagraphStyle(
                'CustomTitle',
                parent=styles['Heading1'],
                fontSize=24,
                textColor='darkblue',
                alignment=TA_CENTER,
                spaceAfter=30
            )
            
            heading_style = ParagraphStyle(
                'CustomHeading',
                parent=styles['Heading2'],
                fontSize=16,
                textColor='darkblue',
                spaceBefore=20,
                spaceAfter=10
            )
            
            subheading_style = ParagraphStyle(
                'CustomSubheading',
                parent=styles['Heading3'],
                fontSize=14,
                textColor='darkgreen',
                spaceBefore=15,
                spaceAfter=8
            )
            
            body_style = ParagraphStyle(
                'CustomBody',
                parent=styles['Normal'],
                fontSize=11,
                alignment=TA_JUSTIFY,
                spaceAfter=10
            )
            
            # Generate images for identified sections
            generated_images = {}
            if 'image_sections' in content_data:
                st.info("🎨 Generating images for your PDF...")
                progress_bar = st.progress(0)
                
                for i, img_section in enumerate(content_data['image_sections']):
                    try:
                        image_bytes = self.image_gen.generate_image(img_section['image_prompt'])
                        if image_bytes:
                            # Save image temporarily
                            img_file = tempfile.NamedTemporaryFile(delete=False, suffix='.png')
                            img_file.write(image_bytes)
                            img_file.close()
                            generated_images[img_section['section_index']] = img_file.name
                        
                        progress_bar.progress((i + 1) / len(content_data['image_sections']))
                    except Exception as e:
                        logger.warning(f"Failed to generate image for section {img_section['title']}: {e}")
                        continue
            
            # Parse content and create PDF elements
            story = []
            content = content_data.get('content', '') if isinstance(content_data, dict) else str(content_data)
            lines = content.split('\n')
            current_section = 0
            
            for line in lines:
                line = line.strip()
                if not line:
                    story.append(Spacer(1, 6))
                    continue
                
                # Check if this line starts a new section that has an image
                if line.startswith('# ') or line.startswith('## '):
                    # Add image if available for this section
                    if current_section in generated_images:
                        try:
                            img_path = generated_images[current_section]
                            img = RLImage(img_path, width=5*inch, height=3*inch)
                            story.append(img)
                            story.append(Spacer(1, 12))
                        except Exception as e:
                            logger.warning(f"Failed to add image to PDF: {e}")
                    current_section += 1
                
                # Add text content
                if line.startswith('# '):
                    story.append(Paragraph(line[2:], title_style))
                elif line.startswith('## '):
                    story.append(Paragraph(line[3:], heading_style))
                elif line.startswith('### '):
                    story.append(Paragraph(line[4:], subheading_style))
                else:
                    if line:
                        story.append(Paragraph(line, body_style))
            
            doc.build(story)
            buffer.seek(0)
            
            # Clean up temporary image files
            for img_path in generated_images.values():
                try:
                    os.unlink(img_path)
                except Exception as e:
                    logger.warning(f"Failed to clean up image file {img_path}: {e}")
            
            return buffer.getvalue()
        
        except Exception as e:
            logger.error(f"PDF generation error: {e}")
            st.error(f"Error creating PDF: {str(e)}")
            return None

class UIComponents:
    """Handles UI component rendering"""
    
    @staticmethod
    def render_sidebar():
        """Render sidebar navigation"""
        st.sidebar.image("https://img.icons8.com/dusk/64/000000/student-center.png", width=64)
        
        # Language Selection
        st.sidebar.markdown("## 🌐 Language / భాష / ಭಾಷೆ / भाषा")
        languages = LanguageManager.SUPPORTED_LANGUAGES
        selected_language = st.sidebar.selectbox(
            "Choose Language:",
            list(languages.keys()),
            index=list(languages.keys()).index(st.session_state.selected_language)
        )

        # Update session state if language changed
        if selected_language != st.session_state.selected_language:
            st.session_state.selected_language = selected_language
            st.session_state.language_code = languages[selected_language]
            st.rerun()

        st.sidebar.markdown("---")
        st.sidebar.markdown("## 📋 Navigation")
        
        # Navigation buttons if analysis is completed
        if st.session_state.analysis_completed:
            UIComponents._render_navigation_buttons()
    
    @staticmethod
    def _render_navigation_buttons():
        """Render navigation buttons in sidebar"""
        st.sidebar.markdown("### 🎯 Actions Available:")
        
        # Current active view
        active_views = {
            'analysis': not any([
                st.session_state.show_doubt_session,
                st.session_state.show_assessment,
                st.session_state.show_video_script,
                st.session_state.show_video_generation,
                st.session_state.show_conclusion,
                st.session_state.show_personalized_pdf
            ]),
            'doubt_session': st.session_state.show_doubt_session,
            'assessment': st.session_state.show_assessment,
            'video_script': st.session_state.show_video_script,
            'video_generation': st.session_state.show_video_generation,
            'conclusion': st.session_state.show_conclusion,
            'personalized_pdf': st.session_state.show_personalized_pdf
        }
        
        # Navigation buttons
        buttons = [
            ("view_analysis", "analysis"),
            ("doubt_session", "doubt_session"),
            ("assessment", "assessment"),
            ("video_script", "video_script"),
            ("video_generation", "video_generation"),
            ("conclusion", "conclusion"),
            ("personalized_pdf", "personalized_pdf")
        ]
        
        for text_key, view_key in buttons:
            button_type = "primary" if active_views[view_key] else "secondary"
            if st.sidebar.button(LanguageManager.get_text(text_key), 
                               use_container_width=True, type=button_type):
                if view_key == "analysis":
                    SessionManager.set_view("")  # Reset all views
                else:
                    SessionManager.set_view(view_key)
                st.rerun()
        
        st.sidebar.markdown("---")
        
        # New Analysis Button
        if st.sidebar.button(LanguageManager.get_text("new_analysis"), 
                           help="Clear current analysis to analyze a new document", 
                           use_container_width=True):
            SessionManager.reset_analysis()
            st.rerun()
    
    @staticmethod
    def render_file_upload():
        """Render file upload section"""
        st.header(LanguageManager.get_text("upload_header"))
        uploaded_file = st.file_uploader(
            "Choose a PDF or PowerPoint file",
            type=SUPPORTED_FILE_TYPES,
            help="Upload academic papers, lecture slides, or study materials"
        )
        
        if uploaded_file is not None:
            # Validate file
            is_valid, message = DocumentProcessor.validate_file(uploaded_file)
            if not is_valid:
                st.error(f"❌ {message}")
                return None
            
            # Display file details
            UIComponents._display_file_info(uploaded_file)
            
            # Process button
            if st.button(LanguageManager.get_text("analysis_button"), type="primary"):
                return UIComponents._process_document(uploaded_file)
        
        return None
    
    @staticmethod
    def _display_file_info(uploaded_file):
        """Display file information"""
        st.subheader("📋 Document Information")
        col1, col2, col3 = st.columns(3)
        
        with col1:
            st.write(f"**📄 Name:** {uploaded_file.name}")
        with col2:
            st.write(f"**📏 Size:** {uploaded_file.size:,} bytes")
        with col3:
            file_type = 'PDF' if 'pdf' in uploaded_file.type else 'PowerPoint'
            st.write(f"**🔧 Type:** {file_type}")
    
    @staticmethod
    def _process_document(uploaded_file):
        """Process uploaded document"""
        with st.spinner("📄 Processing document and generating comprehensive analysis..."):
            # Extract text based on file type
            if uploaded_file.type == "application/pdf":
                text = DocumentProcessor.extract_text_from_pdf(uploaded_file)
                file_type = "PDF"
            else:
                text = DocumentProcessor.extract_text_from_ppt(uploaded_file)
                file_type = "PowerPoint"
            
            if text:
                st.session_state.document_text = text
                
                # Generate analysis
                content_gen = ContentGenerator()
                analysis = content_gen.analyze_document(text, file_type)
                
                if analysis:
                    st.session_state.analysis_result = analysis
                    st.session_state.analysis_completed = True
                    
                    # Generate multi-level video scripts
                    video_scripts = content_gen.generate_multi_level_video_scripts(text)
                    if video_scripts:
                        st.session_state.video_scripts_by_level = video_scripts
                    
                    # Log the interaction to database
                    try:
                        log_ui_interaction(
                            user_email=st.session_state.user_email,
                            document_name=uploaded_file.name,
                            file_type=file_type,
                            file_size=uploaded_file.size,
                            language_used=st.session_state.selected_language
                        )
                    except Exception as e:
                        logger.warning(f"Failed to log interaction: {e}")
                    
                    st.success("✅ Analysis completed successfully!")
                    st.rerun()
                else:
                    st.error("❌ Failed to generate analysis. Please try again.")
            else:
                st.error("❌ Failed to extract text from the document. Please check the file format.")
        
        return True

# Utility functions
def clean_json_string(raw_text: str) -> str:
    """Clean JSON string for parsing"""
    cleaned = raw_text.strip()
    cleaned = re.sub(r"```json|```", "", cleaned)
    cleaned = cleaned.replace(""", '"').replace(""", '"')
    cleaned = cleaned.replace("'", "'").replace("'", "'")
    return cleaned

def display_translated_content(content: str) -> str:
    """Display content with translation if needed"""
    if st.session_state.language_code != "en":
        with st.spinner(f"📄 Translating to {st.session_state.selected_language}..."):
            return LanguageManager.translate_content(content, st.session_state.language_code)
    return content

# View rendering functions
def render_doubt_session():
    """Render doubt session view"""
    st.header("🤔 Ask Your Doubts")
    st.markdown("Feel free to ask any questions about the document content. I'm here to help clarify your doubts!")
    
    student_question = st.text_area(
        "💭 What would you like to understand better?",
        placeholder="Type your question about the document content here...",
        height=100
    )
    
    if st.button("Get Answer", type="primary"):
        if student_question.strip():
            with st.spinner("🤔 Thinking about your question..."):
                content_gen = ContentGenerator()
                answer = content_gen.answer_doubt(student_question, st.session_state.document_text)
                
                if answer:
                    try:
                        update_ui_interaction(st.session_state.user_email, doubt_sessions=1)
                    except Exception as e:
                        logger.warning(f"Failed to update interaction: {e}")
                    
                    st.markdown("### 📚 Professor's Response:")
                    translated_answer = display_translated_content(answer)
                    st.markdown(translated_answer)
                else:
                    st.error("❌ Sorry, I couldn't process your question. Please try again.")
        else:
            st.warning("⚠️ Please enter a question first.")

def render_assessment():
    """Render assessment view"""
    st.header("📝 Knowledge Assessment")
    
    # Generate questions if not already generated
    if st.session_state.assessment_questions is None:
        with st.spinner("📄 Generating personalized questions..."):
            content_gen = ContentGenerator()
            questions_json = content_gen.generate_assessment(st.session_state.document_text)
            
            if questions_json:
                try:
                    cleaned_json = clean_json_string(questions_json)
                    st.session_state.assessment_questions = json.loads(cleaned_json)
                except json.JSONDecodeError as e:
                    logger.error(f"JSON parsing error: {e}")
                    st.error("❌ Error generating questions. Please try again.")
                    st.session_state.assessment_questions = None
    
    # Display questions if available
    if st.session_state.assessment_questions:
        st.markdown("**Instructions:** Choose the best answer for each question.")
        
        # Create form for questions
        with st.form("assessment_form"):
            answers = {}
            
            # Render questions dynamically
            for q_num in ["question1", "question2"]:
                if q_num in st.session_state.assessment_questions:
                    question_data = st.session_state.assessment_questions[q_num]
                    
                    st.markdown(f"### Question {q_num[-1]}")
                    st.markdown(question_data["question"])
                    
                    key = f"q{q_num[-1]}"
                    answers[key] = st.radio(
                        "Select your answer:",
                        list(question_data["options"].keys()),
                        format_func=lambda x, qd=question_data: f"{x}. {qd['options'][x]}",
                        key=key
                    )
                    
                    if q_num != "question2":  # Don't add separator after last question
                        st.markdown("---")
            
            # Submit button
            submitted = st.form_submit_button("📊 Submit Assessment", type="primary")
            
            if submitted:
                _process_assessment_results(answers)

def _process_assessment_results(answers: Dict[str, str]):
    """Process and display assessment results"""
    # Calculate score
    correct_answers = {
        "q1": st.session_state.assessment_questions["question1"]["correct_answer"],
        "q2": st.session_state.assessment_questions["question2"]["correct_answer"]
    }
    
    score = 0
    total = 2
    
    st.markdown("## 📊 Assessment Results")
    
    # Process each question
    for i, (q_key, correct_answer) in enumerate(correct_answers.items(), 1):
        st.markdown(f"### Question {i} Feedback")
        if answers[q_key] == correct_answer:
            st.success("✅ Correct!")
            score += 1
        else:
            st.error(f"❌ Incorrect. The correct answer is {correct_answer}")
        
        question_key = f"question{i}"
        explanation = st.session_state.assessment_questions[question_key]["explanation"]
        st.markdown(f"**Explanation:** {explanation}")
        
        if i < len(correct_answers):
            st.markdown("---")
    
    # Overall score
    percentage = (score / total) * 100
    st.markdown("### 🎯 Overall Performance")
    
    if percentage >= 50:
        st.success(f"🎉 Great job! You scored {score}/{total} ({percentage:.0f}%)")
        performance_level = "Good"
    else:
        st.warning(f"📚 You scored {score}/{total} ({percentage:.0f}%). Keep studying!")
        performance_level = "Needs Improvement"
    
    # Store quiz performance for PDF generation
    st.session_state.quiz_performance = {
        "score": score,
        "total": total,
        "percentage": percentage,
        "level": performance_level
    }
    
    try:
        update_ui_interaction(
            st.session_state.user_email, 
            assessments_taken=1,
            quiz_score=f"{score}/{total} ({percentage:.0f}%)"
        )
    except Exception as e:
        logger.warning(f"Failed to update interaction: {e}")

def render_video_script():
    """Render video script view with multi-level scripts"""
    st.header("🎬 Educational Video Scripts")
    
    # Generate video scripts if not already generated
    if not st.session_state.video_scripts_by_level:
        with st.spinner("🎥 Creating multi-level video scripts..."):
            content_gen = ContentGenerator()
            scripts = content_gen.generate_multi_level_video_scripts(st.session_state.document_text)
            
            if scripts:
                st.session_state.video_scripts_by_level = scripts
                try:
                    update_ui_interaction(st.session_state.user_email, video_scripts_generated=len(scripts))
                except Exception as e:
                    logger.warning(f"Failed to update interaction: {e}")
            else:
                st.error("❌ Failed to generate video scripts. Please try again.")
    
    # Display video scripts
    if st.session_state.video_scripts_by_level:
        st.markdown("### Choose your preferred difficulty level:")
        
        # Create tabs for different levels
        level_tabs = st.tabs([
            LanguageManager.get_text("easy_level"),
            LanguageManager.get_text("medium_level"), 
            LanguageManager.get_text("hard_level")
        ])
        
        script_levels = ["easy", "medium", "hard"]
        
        for i, (tab, level) in enumerate(zip(level_tabs, script_levels)):
            with tab:
                if level in st.session_state.video_scripts_by_level:
                    translated_script = display_translated_content(st.session_state.video_scripts_by_level[level])
                    st.markdown(translated_script)
                    
                    # Download button for each script
                    st.download_button(
                        label=f"📥 Download {level.title()} Script",
                        data=st.session_state.video_scripts_by_level[level],
                        file_name=f"video_script_{level}.md",
                        mime="text/markdown",
                        help=f"Download the {level} level video script",
                        key=f"download_{level}_script"
                    )
                else:
                    st.info(f"Script for {level} level not available.")
        
        # Button to proceed to video generation
        st.markdown("---")
        st.markdown("### 🎥 Ready to create videos?")
        if st.button("🎬 Generate Videos from Scripts", type="primary", help="Create AI-generated videos using Synthesia"):
            SessionManager.set_view("video_generation")
            st.rerun()

def render_video_generation():
    """Render video generation view using Synthesia API"""
    st.header("🎥 AI Video Generation with Synthesia")
    st.markdown("Transform your video scripts into professional AI-generated videos with virtual presenters!")
    
    # Check if we have video scripts
    if not st.session_state.video_scripts_by_level:
        st.warning("⚠️ Please generate video scripts first before creating videos.")
        if st.button("📝 Go to Video Scripts"):
            SessionManager.set_view("video_script")
            st.rerun()
        return
    
    # Initialize Synthesia generator
    synthesia_gen = get_synthesia_generator()
    
    # Check API configuration
    if not SYNTHESIA_API_KEY or SYNTHESIA_API_KEY == "your-synthesia-api-key-here":
        st.error("❌ Synthesia API key not configured. Please add your API key to generate videos.")
        st.markdown("""
        **To configure Synthesia API:**
        1. Sign up at [Synthesia.io](https://synthesia.io)
        2. Get your API key from the dashboard
        3. Add it to the `SYNTHESIA_API_KEY` constant in the code
        """)
        return
    
    # Video generation section
    st.markdown("### 🎬 Select Scripts to Convert to Videos")
    
    available_scripts = st.session_state.video_scripts_by_level
    script_levels = ["easy", "medium", "hard"]
    level_names = {
        "easy": LanguageManager.get_text("easy_level"),
        "medium": LanguageManager.get_text("medium_level"),
        "hard": LanguageManager.get_text("hard_level")
    }
    
    # Create checkboxes for each available script
    selected_levels = []
    for level in script_levels:
        if level in available_scripts:
            if st.checkbox(f"Generate video for {level_names[level]}", key=f"select_{level}"):
                selected_levels.append(level)
    
    if not selected_levels:
        st.info("📺 Select at least one script level to generate videos.")
        return
    
    # Video generation settings
    with st.expander("🔧 Video Generation Settings"):
        st.markdown("**Avatar and Voice Selection:**")
        
        # Show current avatar assignments
        for level in selected_levels:
            st.markdown(f"- **{level_names[level]}:** Avatar will be selected automatically based on difficulty level")
        
        st.markdown("""
        **Video Specifications:**
        - **Format:** MP4
        - **Aspect Ratio:** 16:9 (1920x1080)
        - **Duration:** Approximately 2-4 minutes per script
        - **Background:** Custom color based on difficulty level
        - **Voice:** Natural AI voice matching the difficulty level
        """)
    
    # Generate videos button
    if st.button("🚀 Generate Selected Videos", type="primary"):
        _generate_synthesia_videos(selected_levels, available_scripts, synthesia_gen)
    
    # Display existing video status
    if st.session_state.generated_videos:
        st.markdown("---")
        st.markdown("### 📺 Your Generated Videos")
        _display_video_status(synthesia_gen)

def _generate_synthesia_videos(selected_levels: List[str], scripts: Dict[str, str], synthesia_gen):
    """Generate videos using Synthesia API"""
    st.markdown("### 🎬 Generating Videos...")
    
    progress_bar = st.progress(0)
    status_container = st.container()
    
    for i, level in enumerate(selected_levels):
        with status_container:
            st.markdown(f"**🎥 Generating {level.title()} Level Video...**")
            
            try:
                # Create video using Synthesia API
                video_result = synthesia_gen.create_video(
                    script=scripts[level],
                    difficulty_level=level,
                    title=f"Educational Content - {level.title()} Level"
                )
                
                if video_result:
                    # Store video information
                    if 'generated_videos' not in st.session_state:
                        st.session_state.generated_videos = {}
                    
                    st.session_state.generated_videos[level] = video_result
                    
                    st.success(f"✅ {level.title()} level video creation started!")
                    st.markdown(f"**Video ID:** `{video_result['video_id']}`")
                    st.markdown(f"**Status:** {video_result['status']}")
                    
                    # Log the video generation
                    try:
                        update_ui_interaction(
                            st.session_state.user_email, 
                            videos_generated=1
                        )
                    except Exception as e:
                        logger.warning(f"Failed to update interaction: {e}")
                    
                else:
                    st.error(f"❌ Failed to create {level} level video. Please try again.")
                    
            except Exception as e:
                st.error(f"❌ Error generating {level} video: {str(e)}")
                logger.error(f"Video generation error for {level}: {e}")
        
        # Update progress
        progress_bar.progress((i + 1) / len(selected_levels))
    
    st.success("🎉 Video generation process completed! Videos are being processed by Synthesia.")
    st.info("💡 Video processing typically takes 5-10 minutes. Check the status below for updates.")

def _display_video_status(synthesia_gen):
    """Display status of generated videos"""
    for level, video_info in st.session_state.generated_videos.items():
        with st.expander(f"📹 {level.title()} Level Video Status", expanded=True):
            col1, col2 = st.columns(2)
            
            with col1:
                st.markdown(f"**Video ID:** `{video_info['video_id']}`")
                st.markdown(f"**Title:** {video_info['title']}")
                st.markdown(f"**Created:** {video_info['created_at']}")
            
            with col2:
                # Check current status
                if st.button(f"🔄 Check Status", key=f"status_{level}"):
                    with st.spinner("Checking video status..."):
                        current_status = synthesia_gen.get_video_status(video_info['video_id'])
                        
                        if current_status:
                            # Update stored status
                            st.session_state.generated_videos[level].update(current_status)
                            st.rerun()
                        else:
                            st.error("Failed to check status")
                
                # Display current status
                current_status = video_info.get('status', 'unknown')
                if current_status == 'complete':
                    st.success("✅ Video Ready!")
                elif current_status == 'in_progress':
                    st.info("🔄 Processing...")
                elif current_status == 'failed':
                    st.error("❌ Generation Failed")
                else:
                    st.warning(f"📊 Status: {current_status}")
            
            # Download button if video is ready
            if video_info.get('status') == 'complete':
                download_url = synthesia_gen.get_video_download_url(video_info['video_id'])
                if download_url:
                    st.markdown(f"**[📥 Download Video]({download_url})**")
                else:
                    st.warning("Download URL not available yet. Please check again in a moment.")
            
            # Show video preview if available
            if 'download' in video_info and video_info['download']:
                st.video(video_info['download'])

def render_conclusion():
    """Render conclusion view"""
    st.header("🎯 Conclusion & Applications")
    
    # Generate conclusion if not already generated
    if not st.session_state.conclusion_content:
        with st.spinner("📄 Generating comprehensive conclusion..."):
            content_gen = ContentGenerator()
            conclusion = content_gen.generate_conclusion(st.session_state.document_text)
            
            if conclusion:
                st.session_state.conclusion_content = conclusion
            else:
                st.error("❌ Failed to generate conclusion. Please try again.")
    
    # Display conclusion
    if st.session_state.conclusion_content:
        translated_conclusion = display_translated_content(st.session_state.conclusion_content)
        st.markdown(translated_conclusion)

def render_personalized_pdf():
    """Render personalized PDF view"""
    st.header("📄 Personalized Course PDF with AI-Generated Images")
    st.markdown("Generate a comprehensive study guide with custom illustrations tailored to your learning needs and quiz performance.")
    
    # Show quiz performance if available
    if st.session_state.quiz_performance:
        col1, col2, col3 = st.columns(3)
        with col1:
            st.metric("Quiz Score", f"{st.session_state.quiz_performance['score']}/{st.session_state.quiz_performance['total']}")
        with col2:
            st.metric("Percentage", f"{st.session_state.quiz_performance['percentage']:.0f}%")
        with col3:
            st.metric("Performance Level", st.session_state.quiz_performance['level'])
        
        st.info("📊 Your PDF will be customized based on your quiz performance!")
    else:
        st.info("💡 Take the assessment first for a more personalized PDF experience!")
    
    # Generate PDF button
    if st.button("🎨 Generate Personalized PDF with Images", type="primary"):
        with st.spinner("📚 Creating your personalized study guide with AI-generated images..."):
            content_gen = ContentGenerator()
            pdf_content_data = content_gen.generate_pdf_content(
                st.session_state.document_text, 
                st.session_state.quiz_performance
            )
            
            if pdf_content_data:
                # Create PDF with images
                pdf_gen = PDFGenerator()
                pdf_bytes = pdf_gen.create_pdf_from_content(pdf_content_data)
                
                if pdf_bytes:
                    st.success("✅ Your personalized PDF with AI-generated images has been created!")
                    try:
                        update_ui_interaction(st.session_state.user_email, pdfs_generated=1)
                    except Exception as e:
                        logger.warning(f"Failed to update interaction: {e}")
                    
                    # Show image generation summary
                    if 'image_sections' in pdf_content_data and pdf_content_data['image_sections']:
                        st.info(f"🎨 Generated {len(pdf_content_data['image_sections'])} custom illustrations for your PDF!")
                    
                    # Download button
                    st.download_button(
                        label="📥 Download Personalized Course Guide with Images",
                        data=pdf_bytes,
                        file_name="personalized_course_guide_with_images.pdf",
                        mime="application/pdf",
                        help="Download your customized study guide PDF with AI-generated illustrations"
                    )
                    
                    # Preview content
                    with st.expander("👀 Preview PDF Content"):
                        content = pdf_content_data.get('content', '')
                        st.markdown(content)
                    
                    # Show image prompts used
                    if 'image_sections' in pdf_content_data and pdf_content_data['image_sections']:
                        with st.expander("🎨 AI Image Generation Details"):
                            st.markdown("**Images generated for the following sections:**")
                            for img_section in pdf_content_data['image_sections']:
                                st.markdown(f"- **{img_section['title']}**")
                                st.markdown(f"  - *Image prompt: {img_section['image_prompt']}*")
                else:
                    st.error("❌ Failed to create PDF. Please try again.")
            else:
                st.error("❌ Failed to generate PDF content. Please try again.")

def render_analysis_results():
    """Render analysis results view"""
    st.header("📚 Document Analysis Results")
    translated_analysis = display_translated_content(st.session_state.analysis_result)
    st.markdown(translated_analysis)
    
    # Action buttons at bottom of analysis
    st.markdown("---")
    st.markdown("### 🚀 Next Steps")
    col1, col2, col3 = st.columns(3)
    
    with col1:
        if st.button("Ask Questions", help="Get clarification on any doubts"):
            SessionManager.set_view("doubt_session")
            st.rerun()
    
    with col2:
        if st.button("Take Assessment", help="Test your understanding"):
            SessionManager.set_view("assessment")
            st.rerun()
    
    with col3:
        if st.button("View Video Scripts", help="See educational video scripts"):
            SessionManager.set_view("video_script")
            st.rerun()

def render_main_content():
    """Render main content based on current view"""
    if not st.session_state.analysis_completed:
        # File upload section
        UIComponents.render_file_upload()
        return
    
    # Route to appropriate view
    view_map = {
        'show_doubt_session': render_doubt_session,
        'show_assessment': render_assessment,
        'show_video_script': render_video_script,
        'show_video_generation': render_video_generation,
        'show_conclusion': render_conclusion,
        'show_personalized_pdf': render_personalized_pdf
    }
    
    # Check which view is active
    active_view = None
    for view_key, render_func in view_map.items():
        if st.session_state.get(view_key, False):
            active_view = render_func
            break
    
    # Render active view or default analysis results
    if active_view:
        active_view()
    else:
        render_analysis_results()

def main():
    """Main application function"""
    # Check authentication first
    check_authentication()
    
    # Initialize session state
    SessionManager.initialize()
    
    # Configure Gemini API
    gemini_manager = get_gemini_manager()
    if not gemini_manager.configure():
        st.error("❌ Failed to configure AI service. Please try again later.")
        return
    
    # App header
    st.title("📚 Student Document Analyzer with AI Video Generation")
    if st.session_state.user_email:
        st.markdown(f"Welcome back, **{st.session_state.user_email}**! Upload academic documents to get started with AI-powered analysis and video generation.")
    
    # Render sidebar
    UIComponents.render_sidebar()
    
    # Render main content
    render_main_content()
    
    # Show logout option
    try:
        show_logout_option()
    except Exception as e:
        logger.warning(f"Failed to show logout option: {e}")

if __name__ == "__main__":
    try:
        main()
    except Exception as e:
        logger.error(f"Application error: {e}")
        st.error("❌ An unexpected error occurred. Please refresh the page and try again.")
        st.exception(e)
